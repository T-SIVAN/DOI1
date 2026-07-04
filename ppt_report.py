from __future__ import annotations

from io import BytesIO
from typing import Any, Dict, Iterable, List

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BLUE = RGBColor(68, 114, 196)
LIGHT_BLUE = RGBColor(217, 225, 242)
MID_BLUE = RGBColor(200, 210, 232)
PALE_BLUE = RGBColor(232, 237, 249)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
MUTED = RGBColor(80, 80, 80)


def compact_text(value: Any, max_chars: int) -> str:
    text = " ".join(str(value or "").split())
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "..."


def chunked(items: List[Dict[str, Any]], size: int) -> Iterable[List[Dict[str, Any]]]:
    for start in range(0, len(items), size):
        yield items[start : start + size]


def figure_key(value: Any) -> str:
    return "".join(ch for ch in str(value or "").lower() if ch.isalnum())


def caption_for_figure(analysis: Dict[str, Any], figure_id: Any) -> str:
    target = figure_key(figure_id)
    for item in analysis.get("figure_table_legends") or []:
        current = figure_key(item.get("figure_id"))
        if current and (current == target or current in target or target in current):
            return str(item.get("caption") or "")
    return ""


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, left, top, width, height, fill=WHITE, line=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def set_font(paragraph, size=14, bold=False, color=BLACK):
    paragraph.font.name = "Microsoft YaHei"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def add_textbox(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    *,
    size=16,
    bold=False,
    color=BLACK,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    set_font(p, size=size, bold=bold, color=color)
    return box


def format_cell(cell, text, fill, font_color=BLACK, size=13, bold=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(text or "")
    p.alignment = PP_ALIGN.LEFT
    set_font(p, size=size, bold=bold, color=font_color)


def add_title_slide(prs: Presentation):
    slide = blank_slide(prs)
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)
    add_rect(slide, Inches(0.75), Inches(2.25), Inches(11.85), Inches(1.12), BLUE)
    add_textbox(
        slide,
        "自动化文献与专利结构分析报告",
        Inches(1.0),
        Inches(2.45),
        Inches(11.35),
        Inches(0.5),
        size=28,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        "PDF 自动解析 · 结构化内容表 · 关键图表解读",
        Inches(1.9),
        Inches(3.62),
        Inches(9.55),
        Inches(0.4),
        size=18,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def reference_header(slide, index: int, analysis: Dict[str, Any], suffix: str = ""):
    title = analysis.get("document_title", "Untitled")
    subtitle_parts = []
    if analysis.get("journal_name"):
        subtitle_parts.append(str(analysis.get("journal_name")))
    if analysis.get("doi"):
        subtitle_parts.append(f"DOI: {analysis.get('doi')}")
    if analysis.get("impact_factor"):
        if_note = analysis.get("impact_factor_note")
        suffix_note = f" ({if_note})" if if_note and str(if_note).lower() not in {"none", "待核验"} else ""
        subtitle_parts.append(f"IF: {analysis.get('impact_factor')}{suffix_note}")
    title_text = f"文献{index}. {compact_text(title, 190)}"
    if suffix:
        title_text = f"{title_text} - {suffix}"
    add_textbox(
        slide,
        title_text,
        Inches(0.55),
        Inches(0.12),
        Inches(12.25),
        Inches(0.32),
        size=14,
        bold=True,
        color=BLACK,
    )
    if subtitle_parts:
        add_textbox(
            slide,
            compact_text(" | ".join(subtitle_parts), 210),
            Inches(0.55),
            Inches(0.48),
            Inches(12.25),
            Inches(0.28),
            size=9,
            color=MUTED,
        )


def default_sections(analysis: Dict[str, Any]) -> List[Dict[str, str]]:
    sections = analysis.get("main_content_sections") or []
    if sections:
        return sections[:7]
    framework = analysis.get("writing_framework") or []
    labels = [
        "背景与痛点",
        "核心原理与概念验证",
        "平台/复杂度验证",
        "应用",
        "总结与讨论",
        "实验与分析方法",
    ]
    defaults = []
    for idx, label in enumerate(labels):
        defaults.append(
            {
                "section_name": label,
                "subtopic": "",
                "key_points": framework[idx] if idx < len(framework) else "待根据文献内容补充。",
            }
        )
    return defaults


def add_main_content_slide(prs: Presentation, analysis: Dict[str, Any], index: int):
    slide = blank_slide(prs)
    reference_header(slide, index, analysis)
    add_textbox(
        slide,
        "写作主要内容",
        Inches(0.05),
        Inches(0.94),
        Inches(4.6),
        Inches(0.35),
        size=18,
        bold=True,
    )

    sections = default_sections(analysis)
    rows = len(sections) + 1
    table_shape = slide.shapes.add_table(
        rows,
        3,
        Inches(0.05),
        Inches(1.32),
        Inches(13.0),
        Inches(5.88),
    )
    table = table_shape.table
    table.columns[0].width = Inches(1.55)
    table.columns[1].width = Inches(1.35)
    table.columns[2].width = Inches(10.1)
    table.rows[0].height = Inches(0.95)

    format_cell(table.cell(0, 0), "核心内容", BLUE, WHITE, size=14, bold=True)
    format_cell(table.cell(0, 1), "", BLUE, WHITE, size=14, bold=True)
    format_cell(
        table.cell(0, 2),
        compact_text(analysis.get("main_content", ""), 260),
        BLUE,
        WHITE,
        size=13,
        bold=True,
    )

    for row_idx, section in enumerate(sections, start=1):
        fill = LIGHT_BLUE if row_idx % 2 else MID_BLUE
        table.rows[row_idx].height = Inches(0.72)
        format_cell(table.cell(row_idx, 0), section.get("section_name", ""), fill, size=13)
        format_cell(table.cell(row_idx, 1), section.get("subtopic", ""), fill, size=12)
        format_cell(
            table.cell(row_idx, 2),
            compact_text(section.get("key_points", ""), 210),
            fill,
            size=13,
        )


def add_picture_fit(slide, image_bytes: bytes, left, top, width, height):
    image = Image.open(BytesIO(image_bytes))
    pixel_width, pixel_height = image.size
    if not pixel_width or not pixel_height:
        return

    frame_ratio = width / height
    image_ratio = pixel_width / pixel_height
    if image_ratio >= frame_ratio:
        final_width = width
        final_height = int(width / image_ratio)
    else:
        final_height = height
        final_width = int(height * image_ratio)

    final_left = left + int((width - final_width) / 2)
    final_top = top + int((height - final_height) / 2)
    stream = BytesIO(image_bytes)
    stream.seek(0)
    slide.shapes.add_picture(stream, final_left, final_top, width=final_width, height=final_height)


def add_caption_panel(slide, figures: List[Dict[str, Any]], analysis: Dict[str, Any], left, top, width, height):
    add_textbox(slide, "图例/图注摘录", left, top, width, Inches(0.32), size=16, bold=True)
    panel_top = top + Inches(0.42)
    panel_height = (height - Inches(0.5)) / max(len(figures), 1)

    for row_idx, fig in enumerate(figures):
        item_top = panel_top + row_idx * panel_height
        fill = LIGHT_BLUE if row_idx % 2 == 0 else PALE_BLUE
        add_rect(slide, left, item_top, width, panel_height - Inches(0.08), fill, WHITE)
        fig_id = fig.get("figure_id", f"Fig. {row_idx + 1}")
        caption = caption_for_figure(analysis, fig_id) or fig.get("caption") or fig.get("content_summary") or "未从 PDF 文本层读取到对应图注。"
        add_textbox(
            slide,
            str(fig_id),
            left + Inches(0.12),
            item_top + Inches(0.12),
            Inches(1.05),
            Inches(0.32),
            size=15,
            bold=True,
        )
        add_textbox(
            slide,
            compact_text(caption, 520),
            left + Inches(1.15),
            item_top + Inches(0.1),
            width - Inches(1.35),
            panel_height - Inches(0.25),
            size=12,
        )


def add_figure_slide(
    prs: Presentation,
    analysis: Dict[str, Any],
    index: int,
    figures: List[Dict[str, Any]],
    chunk_index: int,
):
    slide = blank_slide(prs)
    reference_header(slide, index, analysis, f"关键图表解析 {chunk_index}")

    preview_images = analysis.get("preview_images") or []
    image_left = Inches(0.25)
    if preview_images:
        for i, item in enumerate(preview_images[:2]):
            top = Inches(0.95 + i * 3.05)
            label = figures[i].get("figure_id", f"fig{i + 1}") if i < len(figures) else f"fig{i + 1}"
            add_textbox(slide, label, image_left, top - Inches(0.38), Inches(1.0), Inches(0.3), size=15, bold=True)
            add_picture_fit(slide, item["data"], image_left, top, Inches(6.65), Inches(2.6))
    else:
        add_caption_panel(slide, figures, analysis, image_left, Inches(0.95), Inches(6.65), Inches(5.95))

    table_shape = slide.shapes.add_table(
        len(figures) + 1,
        3,
        Inches(7.4),
        Inches(1.28),
        Inches(5.65),
        Inches(5.65),
    )
    table = table_shape.table
    table.columns[0].width = Inches(1.25)
    table.columns[1].width = Inches(1.65)
    table.columns[2].width = Inches(2.75)
    headers = ["主要图示/表", "代表的图示/表格内容", "图示设计目的"]
    for col, header in enumerate(headers):
        format_cell(table.cell(0, col), header, BLUE, WHITE, size=13, bold=True)
    table.rows[0].height = Inches(0.7)

    for row_idx, fig in enumerate(figures, start=1):
        fill = LIGHT_BLUE if row_idx % 2 else PALE_BLUE
        table.rows[row_idx].height = Inches(2.1)
        format_cell(table.cell(row_idx, 0), fig.get("figure_id", ""), fill, size=16)
        format_cell(table.cell(row_idx, 1), compact_text(fig.get("content_summary", ""), 110), fill, size=12)
        format_cell(table.cell(row_idx, 2), compact_text(fig.get("design_purpose", ""), 95), fill, size=12)


def build_pptx_bytes(analyses: List[Dict[str, Any]]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    add_title_slide(prs)

    for index, analysis in enumerate(analyses, start=1):
        add_main_content_slide(prs, analysis, index)
        figures = analysis.get("figures_analysis") or []
        if not figures and analysis.get("figure_table_legends"):
            figures = [
                {
                    "figure_id": item.get("figure_id", f"Fig. {idx}"),
                    "caption": item.get("caption", ""),
                    "content_summary": compact_text(item.get("caption", ""), 110),
                    "design_purpose": "依据 PDF 图注整理，需结合正文进一步确认展示目的。",
                }
                for idx, item in enumerate(analysis.get("figure_table_legends")[:8], start=1)
            ]
        if not figures:
            figures = [
                {
                    "figure_id": "fig1",
                    "content_summary": "模型未提取到明确图表，可尝试结构化模式或 OCR 后重新上传 PDF。",
                    "design_purpose": "提示用户补充可读取的图注文本。",
                }
            ]
        for chunk_index, chunk in enumerate(chunked(figures, 2), start=1):
            add_figure_slide(prs, analysis, index, chunk, chunk_index)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()
