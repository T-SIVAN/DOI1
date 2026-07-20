from __future__ import annotations

import json
from io import BytesIO
from pathlib import Path
from typing import Any, Iterable, List, Tuple

import pandas as pd


SUPPORTED_ATTACHMENT_TYPES = [
    "pdf",
    "docx",
    "pptx",
    "txt",
    "md",
    "csv",
    "tsv",
    "xlsx",
    "xls",
    "json",
    "ris",
    "bib",
]


def _compact(value: Any, limit: int) -> str:
    text = "\n".join(line.rstrip() for line in str(value or "").splitlines())
    text = text.strip()
    if len(text) <= limit:
        return text
    return text[: max(0, limit - 30)].rstrip() + "\n……（附件内容已截断）"


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def extract_attachment_text(file_name: str, data: bytes, max_chars: int = 12_000) -> str:
    """Extract user-visible text from common research attachment formats."""
    suffix = Path(file_name or "attachment").suffix.lower()

    if suffix == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(data))
        pages = [(page.extract_text() or "") for page in reader.pages[:40]]
        return _compact("\n\n".join(pages), max_chars)

    if suffix == ".docx":
        from docx import Document

        document = Document(BytesIO(data))
        parts: List[str] = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                parts.append("\t".join(cell.text.strip() for cell in row.cells))
        return _compact("\n".join(parts), max_chars)

    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(BytesIO(data))
        parts: List[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_text = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if slide_text:
                parts.append(f"[第 {index} 页]\n" + "\n".join(slide_text))
        return _compact("\n\n".join(parts), max_chars)

    if suffix in {".xlsx", ".xls"}:
        workbook = pd.ExcelFile(BytesIO(data))
        parts = []
        for sheet_name in workbook.sheet_names[:8]:
            frame = pd.read_excel(workbook, sheet_name=sheet_name).head(200)
            parts.append(f"[工作表：{sheet_name}]\n{frame.to_csv(index=False)}")
        return _compact("\n\n".join(parts), max_chars)

    if suffix in {".csv", ".tsv"}:
        separator = "\t" if suffix == ".tsv" else ","
        frame = pd.read_csv(BytesIO(data), sep=separator).head(300)
        return _compact(frame.to_csv(index=False), max_chars)

    raw_text = _decode_text(data)
    if suffix == ".json":
        try:
            raw_text = json.dumps(json.loads(raw_text), ensure_ascii=False, indent=2)
        except json.JSONDecodeError:
            pass
    return _compact(raw_text, max_chars)


def build_attachment_context(
    uploaded_files: Iterable[Any] | None,
    max_chars: int = 24_000,
) -> Tuple[str, List[str]]:
    """Return prompt-ready attachment text and readable extraction warnings."""
    sections: List[str] = []
    warnings: List[str] = []
    remaining = max_chars

    for uploaded in uploaded_files or []:
        name = str(getattr(uploaded, "name", "attachment"))
        try:
            data = uploaded.getvalue()
            text = extract_attachment_text(name, data, max_chars=min(12_000, remaining))
            if not text.strip():
                warnings.append(f"{name}：未提取到可读文本（可能是扫描件）。")
                continue
            section = f"### 附件：{name}\n{text}"
            sections.append(section)
            remaining -= len(section)
            if remaining <= 200:
                warnings.append("附件总内容较长，已按模型上下文上限截断。")
                break
        except Exception as exc:
            warnings.append(f"{name}：读取失败（{exc}）")

    return "\n\n".join(sections), warnings
